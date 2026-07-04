"""Tests for POST /api/report/pptx — live-data situation brief as a PPTX deck (§8).

Pure/offline: builds the deck from the posted payload, no upstream HTTP. Re-opens
the returned bytes with python-pptx to prove it's a valid, populated deck.
"""

from __future__ import annotations

import io

from app.routes.export import _BriefKpi, _BriefPayload, report_pptx


def test_report_pptx_returns_valid_populated_deck() -> None:
    payload = _BriefPayload(
        title="Situation brief",
        kpis=_BriefKpi(contacts=17737, feeds_live=6, feeds_total=6, alerts=2),
        severity={"critical": 0, "high": 0, "medium": 0, "low": 2, "info": 0},
        alerts=['LOW: Saved search "all": 37 new'],
        sources=["Aircraft — Global — green", "Vessels — live — green"],
    )
    resp = report_pptx(payload)
    assert resp.status_code == 200
    assert "presentationml" in resp.media_type
    body = resp.body
    assert body[:2] == b"PK", "PPTX is a zip container"

    from pptx import Presentation

    prs = Presentation(io.BytesIO(body))
    slides = list(prs.slides)
    # title + current-picture + recent-alerts + sources
    assert len(slides) == 4
    # the contact count lands in the picture slide text
    all_text = " ".join(
        shape.text_frame.text
        for slide in slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "17,737" in all_text
    assert "Saved search" in all_text


def test_report_pptx_minimal_payload_still_valid() -> None:
    """An empty picture (no alerts/sources) still yields a valid 2-slide deck."""
    resp = report_pptx(_BriefPayload())
    assert resp.status_code == 200
    from pptx import Presentation

    prs = Presentation(io.BytesIO(resp.body))
    assert len(list(prs.slides)) == 2  # title + current picture only
