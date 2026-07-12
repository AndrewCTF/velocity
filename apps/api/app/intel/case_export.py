"""Case → report: the other end of the investigation loop (P2).

Walks a Situation's linked children, their sourced assertions, and attached
evidence (with hashes + custody) into a shareable case report
(docs/roadmap-practitioners-2026-07.md §P2). The value line: every claim in the
output carries a provenance footnote — "asserted by <source> at <observed_at>;
evidence sha256:<hash>" — so a skeptic, an editor, or a court can check it.

(Named ``case_export`` to avoid colliding with ``intel/dossier.py``, which
assembles entity pattern-of-life reads — a different concept.)

Three renderings share one ``build_bundle``:
- JSON bundle (case + custody manifest) for interchange,
- self-contained HTML report with per-claim provenance footnotes,
- PPTX brief retargeted from the live snapshot to the case scope.

Optional AI-drafted narrative is rendered ONLY inside a visibly-labeled block
(``AI_LABEL``); the accepted-AI-use red line (labeled draft + human sign-off) is
enforced here at render time — a caller cannot smuggle unlabeled generated text
into the document through this module.
"""

from __future__ import annotations

import html
import io
import json
import time
from typing import Any

from app.config import Settings, get_settings
from app.intel import evidence as ev
from app.intel.ontology import get_registry
from app.keys import UserCtx

_SITUATION_KIND = "situation"

# The label that MUST wrap any generated narrative in an evidentiary document.
AI_LABEL = (
    "AI-DRAFTED — UNVERIFIED. This text was machine-generated as a drafting "
    "aid; a human must verify every statement against the cited evidence "
    "before relying on it. Not itself evidence."
)

_PPTX_MEDIA = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


class CaseExportError(Exception):
    """The target is not an exportable situation."""


def _now_iso() -> str:
    return time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())


async def build_bundle(
    ctx: UserCtx, sit_id: str, *, settings: Settings | None = None
) -> dict[str, Any]:
    """Machine-readable case bundle: situation + children + assertions + manifest.

    Every child carries its full assertion history (each assertion has a source,
    observed_at, and confidence — the provenance-footnote data). Evidence
    children additionally get a hash-verified custody manifest.
    """
    settings = settings or get_settings()
    reg = get_registry(ctx, settings)
    sit_obj = await reg.get(sit_id)
    if sit_obj is None or (sit_obj.props or {}).get("kind") != _SITUATION_KIND:
        raise CaseExportError(f"{sit_id} is not a situation")

    around = await reg.traverse(sit_id, depth=1)
    children = [o for o in around.objects if o.id != sit_id]

    child_entries: list[dict[str, Any]] = []
    evidence_ids: list[str] = []
    for o in children:
        assertions = await reg.get_assertions(o.id, limit=500)
        if o.id.startswith(f"{ev.EVIDENCE_KIND}:"):
            evidence_ids.append(o.id)
        child_entries.append(
            {
                "object": o.model_dump(),
                "assertions": [a.model_dump() for a in assertions],
            }
        )

    sit_assertions = await reg.get_assertions(sit_id, limit=500)
    manifest = (
        await ev.custody_manifest(ctx, evidence_ids, settings=settings)
        if evidence_ids
        else None
    )
    return {
        "generated_at": _now_iso(),
        "generated_by": ctx.user_id,
        "situation": sit_obj.model_dump(),
        "situation_assertions": [a.model_dump() for a in sit_assertions],
        "children": child_entries,
        "links": [link.model_dump() for link in around.links],
        "evidence_manifest": manifest,
    }


# ── HTML report ───────────────────────────────────────────────────────────────


def _footnote(a: dict[str, Any]) -> str:
    """One provenance footnote for an assertion (never empty: source+observed_at
    are non-null columns, so every rendered claim is attributable)."""
    src = html.escape(str(a.get("source") or "unknown"))
    obs = html.escape(str(a.get("observed_at") or ""))
    conf = a.get("confidence")
    conf_s = f"; confidence {conf}" if conf is not None else ""
    deriv = a.get("derivation")
    deriv_s = f"; {html.escape(str(deriv))}" if deriv else ""
    return f"asserted by {src} at {obs}{conf_s}{deriv_s}"


def _fmt_value(v: Any) -> str:
    if isinstance(v, dict | list):
        return html.escape(json.dumps(v, ensure_ascii=False)[:500])
    return html.escape(str(v))


def render_html(bundle: dict[str, Any], *, narrative: str | None = None) -> str:
    sit = bundle["situation"]
    props = sit.get("props", {})
    name = html.escape(str(props.get("name") or sit["id"]))
    parts: list[str] = []
    parts.append(
        "<!doctype html><html lang=en><head><meta charset=utf-8>"
        f"<title>Case report — {name}</title><style>"
        "body{font:15px/1.55 system-ui,sans-serif;max-width:900px;margin:2rem auto;"
        "padding:0 1rem;color:#111}h1,h2,h3{line-height:1.2}"
        "table{border-collapse:collapse;width:100%;margin:.5rem 0 1.5rem}"
        "th,td{border:1px solid #ccc;padding:.35rem .5rem;text-align:left;"
        "vertical-align:top;font-size:13px}th{background:#f3f4f6}"
        ".fn{color:#555;font-size:11px}"
        ".ev{background:#f9fafb;border-left:3px solid #2563eb;padding:.5rem .75rem;"
        "margin:.5rem 0}"
        ".ai{background:#fff7ed;border:2px dashed #ea580c;padding:.75rem;margin:1rem 0}"
        ".ai b{color:#c2410c}"
        ".meta{color:#666;font-size:12px}code{font-size:12px;word-break:break-all}"
        "</style></head><body>"
    )
    parts.append(f"<h1>Case report: {name}</h1>")
    parts.append(
        f"<p class=meta>Situation <code>{html.escape(sit['id'])}</code> · "
        f"severity {html.escape(str(props.get('severity','')))} · "
        f"status {html.escape(str(props.get('status','')))} · "
        f"generated {html.escape(bundle['generated_at'])} "
        f"by {html.escape(bundle['generated_by'])}</p>"
    )

    summary = str(props.get("summary") or "").strip()
    if summary:
        parts.append("<h2>Analyst summary</h2>")
        parts.append(f"<p>{html.escape(summary)}</p>")
    report = str(props.get("report") or "").strip()
    if report:
        parts.append("<h2>Analyst report</h2>")
        parts.append(f"<p>{html.escape(report).replace(chr(10), '<br>')}</p>")

    if narrative:
        parts.append(
            f"<div class=ai><b>{html.escape(AI_LABEL)}</b>"
            f"<p>{html.escape(narrative).replace(chr(10), '<br>')}</p></div>"
        )

    # Evidence exhibits (hash + custody).
    manifest = bundle.get("evidence_manifest")
    if manifest and manifest.get("items"):
        parts.append("<h2>Evidence exhibits</h2>")
        parts.append(
            f"<p class=meta>{manifest['count']} exhibit(s); "
            f"manifest sha256 <code>{html.escape(manifest['manifest_sha256'])}</code></p>"
        )
        for it in manifest["items"]:
            # "verified" requires a real hash re-check (blob_verified), never
            # mere presence — the report must not vouch for a tampered exhibit.
            if it.get("blob_verified"):
                ok = "verified"
            elif it.get("blob_present"):
                ok = "ALTERED — hash mismatch"
            else:
                ok = "MISSING"
            src_html = (
                f" · source {html.escape(str(it['source_url']))}"
                if it.get("source_url")
                else ""
            )
            parts.append(
                "<div class=ev>"
                f"<b>{html.escape(str(it.get('title') or it['id']))}</b><br>"
                f"<span class=fn>sha256 <code>{html.escape(it['sha256'])}</code> "
                f"({ok}) · {html.escape(str(it.get('media_type') or ''))} · "
                f"{it.get('size_bytes') or 0} bytes · "
                f"method {html.escape(str(it.get('capture_method') or ''))} · "
                f"captured {html.escape(str(it.get('captured_at') or ''))} "
                f"by {html.escape(str(it.get('captured_by') or ''))}"
                f"{src_html}</span></div>"
            )

    # Linked entities and their sourced claims.
    non_ev = [
        c
        for c in bundle["children"]
        if not c["object"]["id"].startswith(f"{ev.EVIDENCE_KIND}:")
    ]
    if non_ev:
        parts.append("<h2>Linked entities &amp; sourced claims</h2>")
    for entry in non_ev:
        obj = entry["object"]
        parts.append(f"<h3><code>{html.escape(obj['id'])}</code></h3>")
        assertions = entry["assertions"]
        if not assertions:
            parts.append(
                "<p class=meta>Linked to the case; no sourced assertions recorded.</p>"
            )
            continue
        parts.append(
            "<table><thead><tr><th>Property</th><th>Value</th>"
            "<th>Provenance</th></tr></thead><tbody>"
        )
        for a in assertions:
            parts.append(
                f"<tr><td>{html.escape(str(a.get('prop')))}</td>"
                f"<td>{_fmt_value(a.get('value'))}</td>"
                f"<td class=fn>{_footnote(a)}</td></tr>"
            )
        parts.append("</tbody></table>")

    parts.append(
        "<hr><p class=meta>Every claim above is footnoted with the source and "
        "time it was asserted; every exhibit is content-addressed by SHA-256. "
        "This document was assembled by Velocity from the local case record.</p>"
    )
    parts.append("</body></html>")
    return "".join(parts)


# ── PPTX brief (case-scoped) ──────────────────────────────────────────────────


def render_pptx(
    bundle: dict[str, Any], *, narrative: str | None = None
) -> bytes | None:
    """Case-scoped PPTX. Returns None if python-pptx is unavailable (caller 503s)."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:  # noqa: BLE001
        return None

    sit = bundle["situation"]
    props = sit.get("props", {})
    name = str(props.get("name") or sit["id"])
    manifest = bundle.get("evidence_manifest") or {}
    non_ev = [
        c
        for c in bundle["children"]
        if not c["object"]["id"].startswith(f"{ev.EVIDENCE_KIND}:")
    ]

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = f"Case report: {name}"
    s.placeholders[1].text = (
        f"severity {props.get('severity','')} · status {props.get('status','')}\n"
        f"Generated {bundle['generated_at']} · keyless OSINT"
    )

    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Case overview"
    tf = s2.placeholders[1].text_frame
    tf.text = f"{len(non_ev)} linked entities"
    for line in (
        f"{manifest.get('count', 0)} evidence exhibits",
        (
            f"Manifest sha256: {manifest.get('manifest_sha256','')[:24]}…"
            if manifest.get("manifest_sha256")
            else ""
        ),
        (str(props.get("summary") or "")[:400]),
    ):
        if not line:
            continue
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)

    if manifest.get("items"):
        s3 = prs.slides.add_slide(prs.slide_layouts[1])
        s3.shapes.title.text = "Evidence exhibits (hashed)"
        etf = s3.placeholders[1].text_frame
        first = manifest["items"][0]
        etf.text = (
            f"{first.get('title') or first['id']} — sha256 {first['sha256'][:16]}…"
        )
        for it in manifest["items"][1:12]:
            p = etf.add_paragraph()
            p.text = f"{it.get('title') or it['id']} — sha256 {it['sha256'][:16]}…"

    if narrative:
        s4 = prs.slides.add_slide(prs.slide_layouts[1])
        s4.shapes.title.text = "Draft narrative"
        ntf = s4.placeholders[1].text_frame
        ntf.text = AI_LABEL
        p = ntf.add_paragraph()
        p.text = narrative[:1200]

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
