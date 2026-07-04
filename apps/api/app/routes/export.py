"""GET /api/export — download the current live picture as GeoJSON or CSV.

Analysts live on "send me that". This reuses the same in-process snapshots the
globe renders (no extra upstream load for aircraft) and emits a portable file:

- ``fmt=geojson`` (default) → a GeoJSON FeatureCollection (RFC 7946), the same
  schema the web adapter consumes — drops straight into QGIS / kepler.gl / Leaflet.
- ``fmt=csv`` → a flat table (id, kind, lon, lat, label, category, course, speed)
  for spreadsheets.

``kinds`` selects layers (``aircraft``, ``vessels``; comma-separated, default
``aircraft``). ``bbox=min_lon,min_lat,max_lon,max_lat`` clips to a viewport.
Vessels are best-effort (keyless Baltic snapshot); an upstream failure yields
zero vessels, never an error — the export of whatever IS live always succeeds.
"""

from __future__ import annotations

import csv
import io
import json
import time
from typing import Any
from xml.sax.saxutils import escape

from fastapi import APIRouter, Query, Response
from pydantic import BaseModel

router = APIRouter(tags=["export"])

_BBox = tuple[float, float, float, float]


def _parse_bbox(raw: str | None) -> _BBox | None:
    if not raw:
        return None
    try:
        parts = [float(x) for x in raw.split(",")]
    except ValueError:
        return None
    if len(parts) != 4:
        return None
    min_lon, min_lat, max_lon, max_lat = parts
    if min_lon >= max_lon or min_lat >= max_lat:
        return None
    return (min_lon, min_lat, max_lon, max_lat)


def _in_bbox(feat: dict[str, Any], bbox: _BBox | None) -> bool:
    if bbox is None:
        return True
    coords = (feat.get("geometry") or {}).get("coordinates")
    if not coords or len(coords) < 2:
        return False
    lon, lat = coords[0], coords[1]
    min_lon, min_lat, max_lon, max_lat = bbox
    return min_lon <= lon <= max_lon and min_lat <= lat <= max_lat


def _tag(feat: dict[str, Any], kind: str) -> dict[str, Any]:
    out = dict(feat)
    out["properties"] = {**(feat.get("properties") or {}), "kind": kind}
    return out


async def _collect(kinds: set[str], bbox: _BBox | None) -> list[dict[str, Any]]:
    # Imported here (not at module load) so the route module never participates
    # in the app-factory import order — matches the codebase's defensive idiom.
    from app.routes import adsb, maritime

    feats: list[dict[str, Any]] = []
    if "aircraft" in kinds:
        snap = await adsb.global_snapshot()
        feats += [_tag(f, "aircraft") for f in snap.get("features", []) if _in_bbox(f, bbox)]
    if "vessels" in kinds:
        try:
            vsnap = await maritime.digitraffic_snapshot()
            feats += [_tag(f, "vessel") for f in vsnap.get("features", []) if _in_bbox(f, bbox)]
        except Exception:
            # Best-effort: vessel upstream down → export aircraft anyway.
            pass
    return feats


def _first(p: dict[str, Any], *keys: str) -> Any:
    """First present (non-None) property among ``keys``, else ""."""
    return next((p[k] for k in keys if p.get(k) is not None), "")


def _to_csv(feats: list[dict[str, Any]]) -> str:
    out = io.StringIO()
    w = csv.writer(out)
    w.writerow(["id", "kind", "lon", "lat", "label", "category", "course", "speed"])
    for f in feats:
        p = f.get("properties") or {}
        coords = (f.get("geometry") or {}).get("coordinates") or []
        w.writerow(
            [
                f.get("id") or p.get("icao24") or p.get("mmsi") or "",
                p.get("kind", ""),
                coords[0] if len(coords) > 0 else "",
                coords[1] if len(coords) > 1 else "",
                p.get("callsign") or p.get("name") or p.get("registration") or "",
                p.get("category") or "",
                _first(p, "track_deg", "cog", "heading"),
                _first(p, "ground_speed", "velocity_ms", "sog"),
            ]
        )
    return out.getvalue()


def _to_kml(feats: list[dict[str, Any]]) -> str:
    parts = [
        '<?xml version="1.0" encoding="UTF-8"?>',
        '<kml xmlns="http://www.opengis.net/kml/2.2"><Document>',
        "<name>Velocity export</name>",
    ]
    for f in feats:
        p = f.get("properties") or {}
        coords = (f.get("geometry") or {}).get("coordinates") or []
        if len(coords) < 2:
            continue
        lon, lat = coords[0], coords[1]
        alt = _first(p, "alt_m", "geo_alt_m", "baro_alt_m") or 0
        label = p.get("callsign") or p.get("name") or p.get("registration") or f.get("id") or ""
        name = escape(str(label))
        desc = escape(str(p.get("kind", "")))
        parts.append(
            f"<Placemark><name>{name}</name><description>{desc}</description>"
            f"<Point><coordinates>{lon},{lat},{alt}</coordinates></Point></Placemark>"
        )
    parts.append("</Document></kml>")
    return "\n".join(parts)


# fmt → (media type, file extension)
_MEDIA = {
    "csv": ("text/csv", "csv"),
    "kml": ("application/vnd.google-earth.kml+xml", "kml"),
    "geojson": ("application/geo+json", "geojson"),
}


@router.get("/api/export")
async def export(
    fmt: str = Query("geojson", pattern="^(geojson|csv|kml)$"),
    kinds: str = Query("aircraft"),
    layer: str | None = Query(None),
    bbox: str | None = Query(None),
    limit: int = Query(0, ge=0, le=50000),
) -> Response:
    # `layer` is the canonical param callers send; `kinds` is the legacy alias.
    # layer wins when present.
    kind_set = {k.strip() for k in (layer or kinds).split(",") if k.strip()} or {"aircraft"}
    feats = await _collect(kind_set, _parse_bbox(bbox))
    if limit:
        feats = feats[:limit]

    media, ext = _MEDIA[fmt]
    if fmt == "csv":
        body = _to_csv(feats)
    elif fmt == "kml":
        body = _to_kml(feats)
    else:
        body = json.dumps({"type": "FeatureCollection", "features": feats})
    return Response(
        content=body,
        media_type=media,
        headers={"Content-Disposition": f'attachment; filename="velocity-export.{ext}"'},
    )


# ── PPTX situation brief (design §8 Slides/Stencil) ─────────────────────────

class _BriefKpi(BaseModel):
    contacts: int = 0
    feeds_live: int = 0
    feeds_total: int = 0
    alerts: int = 0


class _BriefPayload(BaseModel):
    title: str = "Situation brief"
    classification: str = "Unclassified // Open-source intelligence"
    kpis: _BriefKpi = _BriefKpi()
    severity: dict[str, int] = {}
    alerts: list[str] = []
    sources: list[str] = []


_PPTX_MEDIA = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@router.post("/api/report/pptx")
def report_pptx(payload: _BriefPayload) -> Response:
    """Generate a live-data situation brief as a real PPTX deck (design §8). The
    frontend posts the current picture (KPIs/severity/alerts/sources) — the same
    live picture the Brief panel renders — and gets back a .pptx. python-pptx is
    a Phase-1 dep; if the import ever fails we degrade with a 503 rather than 500."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:  # noqa: BLE001 — graceful degrade if the optional lib is missing
        return Response(content="pptx unavailable", status_code=503)

    stamp = time.strftime("%Y-%m-%d %H:%M:%SZ", time.gmtime())
    prs = Presentation()

    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = payload.title
    s.placeholders[1].text = f"{payload.classification}\nGenerated {stamp} · keyless OSINT"

    # Current-picture slide
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Current picture"
    tf = s2.placeholders[1].text_frame
    tf.text = f"{payload.kpis.contacts:,} tracked contacts"
    for line in (
        f"{payload.kpis.feeds_live}/{payload.kpis.feeds_total} feeds live",
        f"{payload.kpis.alerts:,} alerts in buffer",
        "Alerts by severity: "
        + " · ".join(f"{k} {v}" for k, v in payload.severity.items())
        if payload.severity
        else "",
    ):
        if not line:
            continue
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)

    # Recent-alerts slide
    if payload.alerts:
        s3 = prs.slides.add_slide(prs.slide_layouts[1])
        s3.shapes.title.text = "Recent alerts"
        atf = s3.placeholders[1].text_frame
        atf.text = payload.alerts[0]
        for a in payload.alerts[1:12]:
            p = atf.add_paragraph()
            p.text = a

    # Sources slide
    if payload.sources:
        s4 = prs.slides.add_slide(prs.slide_layouts[1])
        s4.shapes.title.text = "Sources"
        stf = s4.placeholders[1].text_frame
        stf.text = payload.sources[0]
        for src in payload.sources[1:20]:
            p = stf.add_paragraph()
            p.text = src

    buf = io.BytesIO()
    prs.save(buf)
    return Response(
        content=buf.getvalue(),
        media_type=_PPTX_MEDIA,
        headers={"Content-Disposition": 'attachment; filename="situation-brief.pptx"'},
    )
